import os
import json
import hashlib

# Paths for JSON storage
USERS_FILE = "users.json"
CURRICULUM_FILE = "curriculum.json"
QNA_FILE = "qna.json"
SLIDES_DIR = "slides"

DEFAULT_CURRICULUM = [
    {
        "id": "1",
        "title": "1회차: 바이브 코딩 입문 및 AI 협업 마인드셋",
        "content": """# 1회차: 바이브 코딩 입문 및 AI 협업 마인드셋

비개발자도 AI 어시스턴트(Antigravity 등)와 함께 소프트웨어를 직접 기획하고 코딩할 수 있는 **'바이브 코딩'**의 기본 개념과 마인드셋을 학습합니다.

---

### 🎯 오늘의 학습 목표
1. **바이브 코딩(Vibe Coding)**의 개념 정의와 전통적 코딩과의 차이점 이해
2. AI 어시스턴트와의 효과적인 협업을 위한 소통 가이드라인 습득
3. 파이썬(Python) 맛보기 및 실행 환경 확인

---

### 💻 가이드 코드 블록
아래 코드는 AI가 비개발자를 위해 작성한 첫 파이썬 스크립트입니다. 
주석(#)을 읽으며 흐름을 파악해 보세요.

```python
# 사용자에게 이름을 물어보고 입력값을 받습니다.
user_name = input("당신의 이름을 입력해주세요: ")

# 환영 메시지를 콘솔에 출력합니다.
print(f"🎉 반갑습니다, {user_name}님! 바이브 코딩의 세계에 오신 것을 환영합니다! 🚀")
```

---

### 💡 실습 미션
- AI 어시스턴트에게 **"사용자로부터 두 개의 숫자를 입력받아 더한 결과를 출력하는 파이썬 코드를 작성해줘"**라고 요청해보세요.
- AI가 작성해 준 코드를 복사하여 로컬 파이썬 파일로 저장하고 직접 실행해보세요!
"""
    },
    {
        "id": "2",
        "title": "2회차: 프롬프트 엔지니어링의 정석 (말하는 대로 만드는 UI)",
        "content": """# 2회차: 프롬프트 엔지니어링의 정석 (말하는 대로 만드는 UI)

원하는 기능과 화면을 명확하게 AI에게 지시하기 위한 **'프롬프트 작성법'**을 마스터합니다. AI가 최적의 결과를 낼 수 있도록 하는 '역할 지정', '제약 조건', '출력 포맷' 설정 방법을 배웁니다.

---

### 🎯 오늘의 학습 목표
1. **구체적 지시(Specific Prompting)**를 통한 오류 최소화
2. **역할 정의(Role-playing)**와 **예시 제공(Few-shot)** 기법의 활용
3. AI에게 간단한 계산기 또는 할 일 목록(Todo) 애플리케이션 기획서 전달해 보기

---

### 💻 가이드 코드 블록
AI에게 화면 구조를 명확히 제안하여 완성된 HTML 코드 구조의 예시입니다.

```html
<!DOCTYPE html>
<html>
<head>
    <title>심플 투두 리스트</title>
    <style>
        body { font-family: sans-serif; padding: 20px; background-color: #f7f9fc; }
        .container { max-width: 400px; margin: auto; background: white; padding: 20px; border-radius: 8px; box-shadow: 0 4px 6px rgba(0,0,0,0.1); }
        .todo-item { display: flex; justify-content: space-between; padding: 8px 0; border-bottom: 1px solid #eee; }
    </style>
</head>
<body>
    <div class="container">
        <h2>📝 오늘의 할 일</h2>
        <input type="text" placeholder="할 일을 입력하세요..." style="width: 70%; padding: 8px;">
        <button style="padding: 8px;">추가</button>
        <div class="todo-item">
            <span>파이썬 문법 복습하기</span>
            <button style="color: red; border: none; background: none; cursor: pointer;">삭제</button>
        </div>
    </div>
</body>
</html>
```

---

### 💡 실습 미션
- AI에게 **"위의 HTML 코드에 JavaScript를 추가하여 실제로 할 일을 입력하고 추가/삭제할 수 있는 동적 기능이 추가된 단일 HTML 파일로 완성해줘"**라고 프롬프트를 작성하여 테스트해 보세요.
"""
    },
    {
        "id": "3",
        "title": "3회차: Streamlit을 활용한 초고속 웹 프로토타입 제작",
        "content": """# 3회차: Streamlit을 활용한 초고속 웹 프로토타입 제작

복잡한 HTML/CSS/JS 없이, 오직 **파이썬 코딩 몇 줄만으로** 아름다운 대시보드와 웹 애플리케이션을 빌드할 수 있는 강력한 프레임워크인 **Streamlit**을 배웁니다.

---

### 🎯 오늘의 학습 목표
1. **Streamlit** 라이브러리의 핵심 동작 원리(단방향 데이터 흐름) 이해
2. `st.write`, `st.button`, `st.text_input`, `st.selectbox` 등 대표 위젯의 사용법 학습
3. 세션 상태(`st.session_state`)를 활용한 데이터 유지 방법 이해

---

### 💻 가이드 코드 블록
Streamlit으로 인터랙티브한 웹 화면을 그리는 간단한 예제 코드입니다.

```python
import streamlit as st

st.title("📊 초간단 인터랙티브 대시보드")

# 사이드바에서 사용자 입력 받기
score = st.sidebar.slider("만족도를 선택하세요:", 0, 100, 80)

# 메인 화면에 출력
st.metric(label="현재 플랫폼 만족도 점수", value=f"{score}점")

if st.button("축하 메시지 보기"):
    st.balloon()
    st.success("교육 과정 참여를 축하합니다!")
```

---

### 💡 실습 미션
- Streamlit을 활용해 **"사용자로부터 몸무게(kg)와 키(cm)를 입력받아 BMI를 계산하고 비만 정도를 알려주는 웹 계산기 앱"**을 작성해 보세요.
"""
    },
    {
        "id": "4",
        "title": "4회차: 공공 API 연동 및 웹 서빙 실습",
        "content": """# 4회차: 공공 API 연동 및 웹 서빙 실습

외부 오픈 API 또는 AI API를 활용하여 실제 데이터가 살아서 움직이는 **'실용 앱'**을 만들고 배포하는 방법을 학습합니다.

---

### 🎯 오늘의 학습 목표
1. **API(Application Programming Interface)**와 JSON 데이터 포맷의 이해
2. `requests` 라이브러리를 활용한 공공 데이터(예: 날씨, 서울시 따릉이 등) 수집
3. Streamlit Cloud를 활용하여 내가 만든 서비스를 웹상에 배포 및 공유하기

---

### 💻 가이드 코드 블록
날씨 API나 목 데이터(Mock Data)를 조회하여 화면에 표 형태로 뿌려주는 샘플입니다.

```python
import streamlit as st
import pandas as pd

st.title("🌦️ 주요 도시 날씨 정보 정보판")

# 가상의 날씨 데이터 생성
weather_data = {
    "도시": ["서울", "부산", "제주", "인천", "대구"],
    "온도 (°C)": [26.5, 24.2, 28.0, 25.1, 27.8],
    "습도 (%)": [65, 70, 80, 60, 55],
    "날씨": ["맑음", "흐림", "비", "맑음", "맑음"]
}

df = pd.DataFrame(weather_data)

# 데이터프레임을 테이블 형태로 깔끔하게 렌더링
st.dataframe(df.style.highlight_max(axis=0, color="#ffe4e1"))
```

---

### 💡 실습 미션
- 공공데이터포털 또는 기타 무료 API 사이트에서 발급받은 API 키를 Streamlit의 Secrets(`.streamlit/secrets.toml`)에 안전하게 설정하고 호출하는 연습을 해보세요.
"""
    }
]

def hash_password(password: str) -> str:
    """Hash a password using SHA-256."""
    return hashlib.sha256(password.encode('utf-8')).hexdigest()

def init_db():
    """Initialize data files if they do not exist."""
    # Ensure slides directory exists
    os.makedirs(SLIDES_DIR, exist_ok=True)
    
    # Initialize users.json
    if not os.path.exists(USERS_FILE):
        with open(USERS_FILE, "w", encoding="utf-8") as f:
            json.dump({}, f, ensure_ascii=False, indent=4)
            
    # Initialize curriculum.json with default 1 to 4 sessions
    if not os.path.exists(CURRICULUM_FILE):
        with open(CURRICULUM_FILE, "w", encoding="utf-8") as f:
            json.dump(DEFAULT_CURRICULUM, f, ensure_ascii=False, indent=4)
            
    # Initialize qna.json
    if not os.path.exists(QNA_FILE):
        with open(QNA_FILE, "w", encoding="utf-8") as f:
            json.dump([], f, ensure_ascii=False, indent=4)

def load_users() -> dict:
    """Load users dict from users.json."""
    init_db()
    try:
        with open(USERS_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}

def save_users(users: dict):
    """Save users dict to users.json."""
    with open(USERS_FILE, "w", encoding="utf-8") as f:
        json.dump(users, f, ensure_ascii=False, indent=4)

def load_curriculum() -> list:
    """Load curriculum list from curriculum.json."""
    init_db()
    try:
        with open(CURRICULUM_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
            # Ensure every item has slides_visible key
            for item in data:
                if "slides_visible" not in item:
                    item["slides_visible"] = True
            return data
    except Exception:
        return []

def save_curriculum(curriculum: list):
    """Save curriculum list to curriculum.json."""
    with open(CURRICULUM_FILE, "w", encoding="utf-8") as f:
        json.dump(curriculum, f, ensure_ascii=False, indent=4)

def load_qna() -> list:
    """Load Q&A list from qna.json.
    Migrates old single-answer format (answer/answered_at) to new multi-answer format (answers:[]).
    """
    init_db()
    try:
        with open(QNA_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
            for item in data:
                # Migrate old format: answer/answered_at -> answers list
                if "answers" not in item:
                    old_answer = item.get("answer")
                    old_at = item.get("answered_at")
                    item["answers"] = []
                    if old_answer:
                        item["answers"].append({
                            "id": "1",
                            "username": "admin",
                            "name": "관리자",
                            "answer": old_answer,
                            "answered_at": old_at or ""
                        })
                    # Clean up old keys
                    item.pop("answer", None)
                    item.pop("answered_at", None)
            return data
    except Exception:
        return []

def save_qna(qna_list: list):
    """Save Q&A list to qna.json."""
    with open(QNA_FILE, "w", encoding="utf-8") as f:
        json.dump(qna_list, f, ensure_ascii=False, indent=4)

def convert_pptx_to_pdf(pptx_path: str, pdf_path: str) -> bool:
    """Convert a PPTX file to PDF using win32com on Windows."""
    try:
        import win32com.client
        import pythoncom
        
        abs_pptx = os.path.abspath(pptx_path)
        abs_pdf = os.path.abspath(pdf_path)
        
        # Initialize COM library for the thread
        pythoncom.CoInitialize()
        
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        # Open PowerPoint presentation without window/gui visible
        deck = powerpoint.Presentations.Open(abs_pptx, WithWindow=False)
        # Save as PDF (Format type 32 is PDF)
        deck.SaveAs(abs_pdf, 32)
        deck.Close()
        powerpoint.Quit()
        return True
    except Exception as e:
        print(f"PPTX to PDF conversion failed: {e}")
        return False

def extract_pptx_notes(pptx_path: str, notes_json_path: str) -> bool:
    """Extract slide notes from a PPTX file and save them as a JSON list."""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        notes_list = []
        for slide in prs.slides:
            slide_notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                slide_notes = slide.notes_slide.notes_text_frame.text.strip()
            notes_list.append(slide_notes)
            
        with open(notes_json_path, "w", encoding="utf-8") as f:
            json.dump(notes_list, f, ensure_ascii=False, indent=4)
        return True
    except Exception as e:
        print(f"Failed to extract PPTX notes: {e}")
        return False

# Initialize databases on import
init_db()
